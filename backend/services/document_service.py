import os
import hashlib
from pathlib import Path
from typing import List, Optional
from pdfplumber import open as pdf_open
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import Chroma
from sqlalchemy.orm import Session
from ..models import File
from .embedding_factory import EmbeddingFactory

# Settings
CHROMA_DIR = os.getenv("CHROMA_DIR", "./chroma_db")
CHUNK_SIZE = 1500
CHUNK_OVERLAP = 200
SEPARATOR = "=================================================="


# file parsing
def parse_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    text = ""

    if ext == ".pdf":
        with pdf_open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text()
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext in [".txt"]:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    return text


# generating hash for files
def file_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


# text splitter
def split_text(text: str) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=[SEPARATOR]
    )
    return splitter.create_documents([text])


# initializing Chroma
def init_chroma_collection(collection_name: str):
    embedding_function = EmbeddingFactory.create()
    return Chroma(
        collection_name=collection_name,
        persist_directory=CHROMA_DIR,
        embedding_function=embedding_function
    )


def init_global_chroma():
    return init_chroma_collection("global_docs")


# save files and processing chunks
def save_file_to_chroma(
        db: Session, user_id: int, file_path: str, filename: str
) -> List[str]:
    """
    upload files, chunking, saving to chroma db
    :param db: database
    :param user_id: user's id
    :param file_path: file path
    :param filename: file name
    :return: chunk_id list
    """
    text = parse_file(file_path)
    f_hash = file_hash(text)

    # check if the file already exists
    existing_file: Optional[File] = db.query(File).filter_by(user_id=user_id, file_hash=f_hash).first()
    if existing_file:
        print("File already processed. Skipping chunking.")
        return []

    # chunking
    chunks = split_text(text)
    print(f"File split into {len(chunks)} chunks")

    # initializing Chroma collection
    collection = init_chroma_collection(f"user_{user_id}")

    # save chunks
    for i, doc in enumerate(chunks):
        collection.add_documents(
            [doc],
            metadatas=[{"file_name": filename, "chunk_index": i, "file_hash": f_hash}],
            ids=[f"{filename}_{i}"]
        )
    collection.persist()
    print(f"Saved {len(chunks)} chunks to Chroma for user {user_id}")

    # save record to database
    new_file = File(user_id=user_id, filename=filename, file_hash=f_hash)
    db.add(new_file)
    db.commit()
    db.refresh(new_file)

    return [f"{filename}_{i}" for i in range(len(chunks))]


# get user's chunks
def get_user_chunks(user_id: int) -> Chroma:
    return init_chroma_collection(f"user_{user_id}")
